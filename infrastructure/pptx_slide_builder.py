import io
import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from domain.i_slide_builder import ISlideBuilder
from config.settings import (
    DEFAULT_TARGET_WIDTH,
    JPEG_QUALITY,
    DPI,
    SLIDE_IMAGE_LEFT_INCH,
    SLIDE_IMAGE_TOP_INCH,
    SLIDE_TEXT_LEFT_INCH,
    SLIDE_TEXT_TOP_INCH,
    SLIDE_TEXT_WIDTH_INCH,
    SLIDE_TEXT_HEIGHT_INCH,
    SLIDE_FONT_SIZE_TIMESTAMP,
    SLIDE_FONT_SIZE_TEXT,
)
from utils.image_utils import resize_image_to_max_width
from utils.time_utils import format_timestamp


class PptxSlideBuilder(ISlideBuilder):
    def __init__(self, output_path: str, video_title: str, video_url: str = ""):
        self.output_path = output_path
        self.video_title = video_title
        self.video_url = video_url
        self.prs = Presentation()
        self.current_slide = None
        self.slide_count = 0

    def create_new_slide(self) -> None:
        slide_layout = self.prs.slide_layouts[6]  # boş layout
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        self.slide_count += 1

    def add_image(
        self,
        image: np.ndarray,
        max_width: int = DEFAULT_TARGET_WIDTH,
        quality: int = JPEG_QUALITY,
        vertical_align: str = "top"
    ) -> None:
        if self.current_slide is None:
            self.create_new_slide()

        resized_bgr = resize_image_to_max_width(image, max_width)
        rgb_image = cv2.cvtColor(resized_bgr, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(rgb_image)

        slide_width_emu = self.prs.slide_width
        slide_height_emu = self.prs.slide_height
        emu_per_inch = 914400
        slide_width_inch = slide_width_emu / emu_per_inch
        slide_height_inch = slide_height_emu / emu_per_inch

        image_width_inch = pil_img.width / DPI
        image_height_inch = pil_img.height / DPI

        if image_width_inch > slide_width_inch:
            scale = slide_width_inch / image_width_inch
            new_width = int(pil_img.width * scale)
            new_height = int(pil_img.height * scale)
            pil_img = pil_img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            image_width_inch = new_width / DPI
            image_height_inch = new_height / DPI

        left_inch = (slide_width_inch - image_width_inch) / 2.0

        if vertical_align == "center":
            top_inch = (slide_height_inch - image_height_inch) / 2.0
        elif vertical_align == "bottom":
            top_inch = slide_height_inch - image_height_inch - 1.0
            if top_inch < 0:
                top_inch = 0
        else:
            top_inch = SLIDE_IMAGE_TOP_INCH

        img_buffer = io.BytesIO()
        pil_img.save(img_buffer, format='JPEG', quality=quality)
        img_buffer.seek(0)

        left = Inches(left_inch)
        top = Inches(top_inch)
        self.current_slide.shapes.add_picture(
            img_buffer, left, top,
            width=Inches(image_width_inch),
            height=Inches(image_height_inch)
        )

    def add_text(self, text: str, font_size: int = SLIDE_FONT_SIZE_TEXT, bold: bool = False) -> None:
        if self.current_slide is None:
            self.create_new_slide()

        left = Inches(SLIDE_TEXT_LEFT_INCH)
        top = Inches(SLIDE_TEXT_TOP_INCH)
        width = Inches(SLIDE_TEXT_WIDTH_INCH)
        height = Inches(SLIDE_TEXT_HEIGHT_INCH)

        textbox = self.current_slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold

    def add_timestamp(self, seconds: float, video_title: str = "") -> None:
        timestamp_str = format_timestamp(seconds, include_ms=True)   # milisaniyeli
        title_str = video_title if video_title else self.video_title
        footer_text = f"{title_str} - {timestamp_str}"
        self.add_text(footer_text, font_size=SLIDE_FONT_SIZE_TIMESTAMP, bold=False)

    def build(self) -> str:
        self.prs.save(self.output_path)
        return self.output_path