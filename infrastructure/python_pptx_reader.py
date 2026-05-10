import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from typing import List, Optional
from domain.i_pptx_reader import IPPTXReader

class PythonPPTXReader(IPPTXReader):
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.prs = Presentation(file_path)

    def get_slide_count(self) -> int:
        return len(self.prs.slides)

    def get_slide_images(self, slide_index: int) -> List[Image.Image]:
        images = []
        slide = self.prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                img = Image.open(io.BytesIO(blob))
                images.append(img)
        return images

    def get_first_image_thumbnail(self, slide_index: int, max_size: int = 100) -> Optional[Image.Image]:
        images = self.get_slide_images(slide_index)
        if not images:
            return None
        img = images[0]
        img.thumbnail((max_size, max_size))
        return img

    def update_from_writer(self, writer) -> None:
        if hasattr(writer, '_new_prs') and writer._new_prs:
            self.prs = writer._new_prs
        elif hasattr(writer, 'last_saved_path') and writer.last_saved_path:
            self.prs = Presentation(writer.last_saved_path)

    def get_slide_text(self, slide_index: int) -> str:
        slide = self.prs.slides[slide_index]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
        return "\n".join(texts)

    def close(self):
        pass