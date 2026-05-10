import io
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from typing import List
from domain.i_pptx_writer import IPPTXWriter

class PythonPPTXWriter(IPPTXWriter):
    def __init__(self, source_path: str):
        self.source_path = source_path
        self._new_prs = None
        self.last_saved_path = None

    def rebuild_with_grid(self, selected_indices: List[int], grid_image: Image.Image) -> None:
        prs = Presentation(self.source_path)
        new_prs = Presentation()
        blank_layout = new_prs.slide_layouts[6]

        for idx, slide in enumerate(prs.slides):
            if idx in selected_indices:
                continue
            new_slide = new_prs.slides.add_slide(blank_layout)
            self._copy_slide_content(slide, new_slide)

        insert_pos = min(selected_indices) if selected_indices else 0
        slides_list = list(new_prs.slides)
        final_prs = Presentation()
        final_layout = final_prs.slide_layouts[6]

        for i in range(insert_pos):
            if i < len(slides_list):
                self._copy_slide_content(slides_list[i], final_prs.slides.add_slide(final_layout))

        grid_slide = final_prs.slides.add_slide(final_layout)
        img_buffer = io.BytesIO()
        grid_image.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(6)
        grid_slide.shapes.add_picture(img_buffer, left, top, width, height)

        for i in range(insert_pos, len(slides_list)):
            self._copy_slide_content(slides_list[i], final_prs.slides.add_slide(final_layout))

        self._new_prs = final_prs

    def delete_slides(self, selected_indices: List[int]) -> None:
        prs = Presentation(self.source_path)
        new_prs = Presentation()
        blank_layout = new_prs.slide_layouts[6]

        for idx, slide in enumerate(prs.slides):
            if idx in selected_indices:
                continue
            new_slide = new_prs.slides.add_slide(blank_layout)
            self._copy_slide_content(slide, new_slide)

        self._new_prs = new_prs

    def replace_slide_image(self, slide_index: int, new_image: Image.Image) -> None:
        prs = self._new_prs if self._new_prs else Presentation(self.source_path)
        slide = prs.slides[slide_index]
        picture_shape = None
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                picture_shape = shape
                break
        if picture_shape is None:
            raise ValueError(f"Slayt {slide_index} içinde resim bulunamadı")
        left, top, width, height = picture_shape.left, picture_shape.top, picture_shape.width, picture_shape.height
        picture_shape.element.getparent().remove(picture_shape.element)
        img_buffer = io.BytesIO()
        new_image.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        slide.shapes.add_picture(img_buffer, left, top, width, height)
        self._new_prs = prs

    def save(self, path: str = None) -> None:
        if path is None:
            path = self.source_path
        self.last_saved_path = path
        if self._new_prs is None:
            prs = Presentation(self.source_path)
            prs.save(path)
        else:
            self._new_prs.save(path)

    def _copy_slide_content(self, src_slide, dst_slide):
        for shape in src_slide.shapes:
            if shape.has_text_frame:
                textbox = dst_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                textbox.text = shape.text
            elif shape.shape_type == 13:  # picture
                blob = shape.image.blob
                dst_slide.shapes.add_picture(io.BytesIO(blob), shape.left, shape.top, shape.width, shape.height)

    def close(self):
        pass